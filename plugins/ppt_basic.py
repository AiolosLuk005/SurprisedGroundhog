# -*- coding: utf-8 -*-
from __future__ import annotations
from pathlib import Path
from core.plugin_base import ExtractResult, register
from core.chunking import Chunk

class PptBasic:
    name = "ppt-basic"
    version = "0.1.0"
    priority = 65

    def can_handle(self, path: str) -> bool:
        return Path(path).suffix.lower() in {'.pptx', '.ppt'}

    def extract(self, path: str, max_chars: int = 4000) -> ExtractResult:
        text = ''
        chunks = []
        p = Path(path)
        try:
            if p.suffix.lower() == '.pptx':
                from pptx import Presentation
                prs = Presentation(path)
                parts = []
                for i, slide in enumerate(prs.slides[:50]):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text.append(shape.text)
                    if slide_text:
                        stxt = "\n".join(slide_text)[:max_chars]
                        chunks.append(
                            Chunk(
                                id=f"{path}#s{i+1}",
                                doc_id=path,
                                text=stxt,
                                page=i+1,
                                metadata={'handler': self.name},
                            )
                        )
                        parts.append(stxt)
                    if sum(len(x) for x in parts) >= max_chars:
                        break
                text = "\n".join(parts)[:max_chars]
            else:
                try:
                    import tempfile, win32com.client
                    tmpdir = tempfile.mkdtemp(prefix="ppt2pptx_plug_")
                    out = str(Path(tmpdir) / (p.stem + ".pptx"))
                    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                    powerpoint.Visible = 0
                    pres = powerpoint.Presentations.Open(str(p), WithWindow=False)
                    pres.SaveAs(out, 24)
                    pres.Close()
                    powerpoint.Quit()
                    from pptx import Presentation
                    prs = Presentation(out)
                    parts = []
                    for i, slide in enumerate(prs.slides[:50]):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slide_text.append(shape.text)
                        if slide_text:
                            stxt = "\n".join(slide_text)[:max_chars]
                            chunks.append(
                                Chunk(
                                    id=f"{path}#s{i+1}",
                                    doc_id=path,
                                    text=stxt,
                                    page=i+1,
                                    metadata={'handler': self.name},
                                )
                            )
                            parts.append(stxt)
                    text = "\n".join(parts)[:max_chars]
                except Exception:
                    text = ''
        except Exception:
            text = ''
        return ExtractResult(text=text, meta={'handler': self.name}, chunks=chunks)

register(PptBasic())
