from __future__ import annotations
import os, io, json, re, urllib.request
from pathlib import Path
from typing import List, Optional

# ---------- 文本抽取 ----------
def extract_text_for_keywords(path: str, max_chars: int = 3000) -> str:
    """
    从多种文档类型中抽取可读文本（前 max_chars 字符）。
    支持：txt/md/rtf/docx/pdf/csv/xlsx/pptx。扫描件 PDF/图片暂不做 OCR。
    """
    p = Path(path)
    ext = p.suffix.lower()
    try:
        if ext in {".txt", ".md", ".rtf", ".log"}:
            with open(p, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()[:max_chars]
        if ext == ".docx":
            from docx import Document
            doc = Document(str(p))
            text = "\n".join(par.text for par in doc.paragraphs if par.text)
            return text[:max_chars]
        if ext == ".pdf":
            from PyPDF2 import PdfReader
            text_parts = []
            reader = PdfReader(str(p))
            for i, page in enumerate(reader.pages[:6]):  # 取前6页
                try:
                    text_parts.append(page.extract_text() or "")
                except Exception:
                    continue
            return "\n".join(text_parts)[:max_chars]
        if ext in {".csv", ".tsv", ".xlsx"}:
            import pandas as pd
            if ext == ".xlsx":
                df = pd.read_excel(str(p), nrows=20)
            else:
                df = pd.read_csv(str(p), nrows=20, sep=None, engine="python")
            text = "\n".join(" ".join(map(str, row)) for row in df.astype(str).values.tolist())
            return text[:max_chars]
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(p))
            texts = []
            for i, slide in enumerate(prs.slides):
                for shp in slide.shapes:
                    if hasattr(shp, "text") and shp.text:
                        texts.append(shp.text)
                if len("".join(texts)) >= max_chars:
                    break
            return "\n".join(texts)[:max_chars]
    except Exception:
        pass
    # 其他类型或失败：返回空串
    return ""

# ---------- 轻量/统计候选 ----------
def _yake_candidates(text: str, lang: str = "zh", topk: int = 12) -> List[str]:
    import yake
    kw = yake.KeywordExtractor(lan=lang, n=1, top=topk)
    pairs = kw.extract_keywords(text or "")
    cand = [w for w, _ in sorted(pairs, key=lambda x: x[1])]  # 分数越小越好
    return _uniq_nonempty(cand)[:topk]

def _jieba_candidates(text: str, topk: int = 12) -> List[str]:
    try:
        import jieba.analyse as ja
        c = ja.extract_tags(text or "", topK=topk)
        return _uniq_nonempty(c)[:topk]
    except Exception:
        return []

def kw_fast(text: str, lang: str = "zh", topk: int = 12) -> List[str]:
    """
    基于 YAKE + jieba 的轻量候选（去重合并）。
    """
    c1 = _yake_candidates(text, lang=lang, topk=topk)
    c2 = _jieba_candidates(text, topk=topk)
    merged = _uniq_nonempty(c1 + c2)
    return merged[:topk]

# ---------- KeyBERT 语义重排 ----------
_KBERT = None
def _get_keybert():
    global _KBERT
    if _KBERT is None:
        from keybert import KeyBERT
        # 轻量多语种模型；若需要中文更强可用 BAAI/bge-small-zh-v1.5
        _KBERT = KeyBERT(model="paraphrase-multilingual-MiniLM-L12-v2")
    return _KBERT

def kw_embed(text: str, candidates: List[str], topk: int = 8) -> List[str]:
    """
    用 KeyBERT 对候选进行语义排序，返回 topk。
    若模型加载失败则回退返回原 candidates。
    """
    try:
        if not text or not candidates:
            return candidates[:topk]
        kb = _get_keybert()
        # KeyBERT 返回 (kw, score) 列表
        scored = kb.extract_keywords(
            text, keyphrase_ngram_range=(1, 3),
            stop_words=None, use_maxsum=True, nr_candidates=min(20, len(candidates)),
            top_n=min(topk, len(candidates)), candidates=candidates
        )
        ranked = [w for w, _ in scored]
        # 合并漏网词，保持稳定
        for w in candidates:
            if w not in ranked and len(ranked) < topk:
                ranked.append(w)
        return _uniq_nonempty(ranked)[:topk]
    except Exception:
        return candidates[:topk]

# ---------- LLM 生成 ----------
def kw_llm(title: str, text: str, seeds: str, max_chars: int = 50,
           model: str = "phi3:mini", host: str = "http://127.0.0.1:11434", timeout: int = 30) -> str:
    """
    通过 Ollama /api/generate 生成中文关键词（逗号分隔，≤max_chars），保留 seeds 前缀。
    """
    prompt = (
        "你是中文关键词提取助手。根据“标题+正文节选”输出中文关键词，要求："
        "1) 只输出关键词，逗号分隔；2) 总长度<={max_len}字；"
        "3) 如果提供了“用户指定关键词(seeds)”，必须把 seeds 放在最前（保持原顺序），再补充其他概括性关键词；"
        "4) 不要解释。\n"
        f"seeds: {seeds or '(无)'}\n"
        f"标题: {title[:80]}\n"
        f"正文节选: {(text or '')[:800]}\n"
        "输出："
    ).replace("{max_len}", str(max_chars))
    data = json.dumps({"model": model, "prompt": prompt, "stream": False}).encode("utf-8")
    req = urllib.request.Request(f"{host}/api/generate", data=data, headers={"Content-Type": "application/json"})
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            x = json.loads(resp.read().decode("utf-8"))
        out = (x.get("response") or "").strip().strip(",， \n")
        # 强制 seeds 前缀
        if seeds:
            out = _force_seeds_prefix(out, seeds)
        return _clip_len(_normalize_commas(out), max_chars)
    except Exception:
        return ""

# ---------- 统一后处理/工具 ----------
def compose_keywords(seeds: str, parts: List[str], max_chars: int = 50) -> str:
    base = _normalize_commas(", ".join(parts))
    if seeds:
        base = _force_seeds_prefix(base, seeds)
    return _clip_len(base, max_chars)

def _normalize_commas(s: str) -> str:
    s = (s or "").replace("，", ",").replace("；", ",").replace(";", ",")
    # 去重空白、重复逗号
    items = [re.sub(r"\s+", "", t) for t in s.split(",")]
    items = [t for t in items if t]
    # 最终使用中文逗号输出
    return "，".join(_uniq_nonempty(items))

def _force_seeds_prefix(out: str, seeds: str) -> str:
    seeds_norm = _normalize_commas(seeds)
    out_norm = _normalize_commas(out)
    # 移除 out 中重复的 seeds 片段，再前置 seeds
    seed_set = set(seeds_norm.split("，"))
    rest = [w for w in out_norm.split("，") if w and w not in seed_set]
    return "，".join(list(seed_set) + rest)

def _clip_len(s: str, max_chars: int) -> str:
    return (s or "")[:max_chars]

def _uniq_nonempty(seq: List[str]) -> List[str]:
    seen, out = set(), []
    for w in seq:
        w = (w or "").strip()
        if not w or w in seen:
            continue
        seen.add(w); out.append(w)
    return out

